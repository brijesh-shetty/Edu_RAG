"""
file_parser.py — Extracts text from PDF, DOCX, and PPTX files.
Returns a list of document dicts: { text, source, page/slide }
"""

import os
import hashlib
import logging
from PIL import Image
import fitz  # PyMuPDF
import pytesseract
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from unstructured.partition.docx import partition_docx
from unstructured.partition.pptx import partition_pptx
from unstructured.chunking.title import chunk_by_title

from src.vision_utils import summarize_image
from src.config import tesseract_path

logger = logging.getLogger(__name__)

# Configure Tesseract Path (for Windows)
try:
    pytesseract.pytesseract.tesseract_cmd = tesseract_path()
except Exception:
    pass

DEFAULT_IMAGES_DIR = os.path.join(os.path.dirname(__file__), "..", "data", "images")
MIN_IMAGE_SIZE_BYTES = 5 * 1024  # 5 KB — skip tiny images (logos, icons, bullets)


def _extract_tables_as_markdown(page) -> str:
    """Extract tables from a PyMuPDF page as markdown using built-in find_tables."""
    try:
        tables = page.find_tables()
        if not tables.tables:
            return ""
        parts = []
        for tab in tables:
            md = tab.to_markdown()
            if md and md.strip():
                parts.append(md)
        if parts:
            return "\n\n" + "\n\n".join(parts)
    except Exception as e:
        logger.debug("Table extraction failed on page: %s", e)
    return ""


def parse_pdf(file_path: str, images_dir: str = None) -> list[dict]:
    """Extract text page-by-page from a PDF file using pymupdf4llm for structured markdown.
    Falls back to OCR for scanned pages. Extracts images with LLaVA summaries."""
    if images_dir is None:
        images_dir = DEFAULT_IMAGES_DIR

    docs = []
    filename = os.path.basename(file_path)
    seen_image_hashes = set()  # Track duplicates across pages

    # Try structured markdown extraction with pymupdf4llm
    try:
        import pymupdf4llm
        page_chunks = pymupdf4llm.to_markdown(file_path, page_chunks=True)
    except Exception as e:
        logger.warning("pymupdf4llm failed, falling back to basic extraction: %s", e)
        page_chunks = None

    pdf = fitz.open(file_path)
    os.makedirs(images_dir, exist_ok=True)
    safe_filename = "".join([c if c.isalnum() else "_" for c in filename])

    for page_num in range(len(pdf)):
        page = pdf[page_num]

        # Use pymupdf4llm markdown if available, else basic extraction
        if page_chunks and page_num < len(page_chunks):
            text = page_chunks[page_num].get("text", "").strip()
        else:
            text = page.get_text().strip()

        # If less than 50 chars, assume scanned page — use OCR
        if len(text) < 50:
            pix = page.get_pixmap(dpi=300)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

            try:
                ocr_text = pytesseract.image_to_string(img)
                text = f"{text}\n{ocr_text}".strip()
            except Exception as e:
                logger.warning("OCR failed on %s page %d: %s", filename, page_num + 1, e)

        # Table extraction fallback/supplement
        table_md = _extract_tables_as_markdown(page)
        if table_md:
            text += table_md

        # Extract images from PDF and summarize with vision model
        try:
            for img_index, img_info in enumerate(page.get_images(full=True)):
                xref = img_info[0]
                base_image = pdf.extract_image(xref)
                image_bytes = base_image["image"]

                # Skip tiny images (logos, icons, bullets)
                if len(image_bytes) < MIN_IMAGE_SIZE_BYTES:
                    logger.debug("Skipping tiny image (%d bytes) on page %d", len(image_bytes), page_num + 1)
                    continue

                # Skip duplicate images (same image on multiple pages)
                img_hash = hashlib.md5(image_bytes).hexdigest()
                if img_hash in seen_image_hashes:
                    logger.debug("Skipping duplicate image on page %d", page_num + 1)
                    continue
                seen_image_hashes.add(img_hash)

                image_ext = base_image["ext"]
                image_filename = f"{safe_filename}_page{page_num+1}_img{img_index}.{image_ext}"
                image_path = os.path.join(images_dir, image_filename)
                with open(image_path, "wb") as f:
                    f.write(image_bytes)

                summary = summarize_image(image_bytes)
                if summary:
                    text += summary
        except Exception as e:
            logger.warning("PDF Image extraction failed: %s", e)

        if text:
            docs.append({
                "text": text,
                "source": filename,
                "page": page_num + 1,
                "type": "pdf"
            })

    pdf.close()
    return docs


def parse_pdf_streaming(file_path: str, images_dir: str = None):
    """Generator that yields (doc_dict, current_page, total_pages) per page for progress tracking."""
    if images_dir is None:
        images_dir = DEFAULT_IMAGES_DIR

    filename = os.path.basename(file_path)

    try:
        import pymupdf4llm
        page_chunks = pymupdf4llm.to_markdown(file_path, page_chunks=True)
    except Exception as e:
        logger.warning("pymupdf4llm failed, falling back to basic extraction: %s", e)
        page_chunks = None

    pdf = fitz.open(file_path)
    total_pages = len(pdf)
    os.makedirs(images_dir, exist_ok=True)
    safe_filename = "".join([c if c.isalnum() else "_" for c in filename])
    seen_image_hashes = set()  # Track duplicates across pages

    for page_num in range(total_pages):
        page = pdf[page_num]

        if page_chunks and page_num < len(page_chunks):
            text = page_chunks[page_num].get("text", "").strip()
        else:
            text = page.get_text().strip()

        if len(text) < 50:
            pix = page.get_pixmap(dpi=300)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            try:
                ocr_text = pytesseract.image_to_string(img)
                text = f"{text}\n{ocr_text}".strip()
            except Exception as e:
                logger.warning("OCR failed on %s page %d: %s", filename, page_num + 1, e)

        table_md = _extract_tables_as_markdown(page)
        if table_md:
            text += table_md

        try:
            for img_index, img_info in enumerate(page.get_images(full=True)):
                xref = img_info[0]
                base_image = pdf.extract_image(xref)
                image_bytes = base_image["image"]

                # Skip tiny images (logos, icons, bullets)
                if len(image_bytes) < MIN_IMAGE_SIZE_BYTES:
                    logger.debug("Skipping tiny image (%d bytes) on page %d", len(image_bytes), page_num + 1)
                    continue

                # Skip duplicate images (same image on multiple pages)
                img_hash = hashlib.md5(image_bytes).hexdigest()
                if img_hash in seen_image_hashes:
                    logger.debug("Skipping duplicate image on page %d", page_num + 1)
                    continue
                seen_image_hashes.add(img_hash)

                image_ext = base_image["ext"]
                image_filename = f"{safe_filename}_page{page_num+1}_img{img_index}.{image_ext}"
                image_path = os.path.join(images_dir, image_filename)
                with open(image_path, "wb") as f:
                    f.write(image_bytes)

                summary = summarize_image(image_bytes)
                if summary:
                    text += summary
        except Exception as e:
            logger.warning("PDF Image extraction failed: %s", e)

        if text:
            doc_dict = {
                "text": text,
                "source": filename,
                "page": page_num + 1,
                "type": "pdf"
            }
            yield doc_dict, page_num + 1, total_pages

    pdf.close()


def parse_docx(file_path: str, images_dir: str = None) -> list[dict]:
    """Extract text semantically from a DOCX file using Unstructured."""
    filename = os.path.basename(file_path)
    docs = []

    try:
        elements = partition_docx(filename=file_path)
        chunks = chunk_by_title(elements)

        for i, chunk in enumerate(chunks, 1):
            if chunk.text.strip():
                docs.append({
                    "text": chunk.text,
                    "source": filename,
                    "page": i,
                    "type": "docx"
                })
    except Exception as e:
        logger.error("Failed to parse docx %s: %s", filename, e)

    return docs


def parse_pptx(file_path: str, images_dir: str = None) -> list[dict]:
    """Extract text semantically from a PPTX file and summarize images."""
    if images_dir is None:
        images_dir = DEFAULT_IMAGES_DIR

    filename = os.path.basename(file_path)
    docs = []

    # 1. Extract image summaries per slide
    slide_summaries = {}
    seen_image_hashes = set()  # Track duplicates across slides
    try:
        prs = Presentation(file_path)

        os.makedirs(images_dir, exist_ok=True)
        safe_filename = "".join([c if c.isalnum() else "_" for c in filename])

        for i, slide in enumerate(prs.slides, 1):
            slide_summaries[i] = ""
            img_index = 0
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_bytes = shape.image.blob

                    # Skip tiny images (logos, icons, bullets)
                    if len(image_bytes) < MIN_IMAGE_SIZE_BYTES:
                        logger.debug("Skipping tiny image (%d bytes) on slide %d", len(image_bytes), i)
                        continue

                    # Skip duplicate images (same image on multiple slides)
                    img_hash = hashlib.md5(image_bytes).hexdigest()
                    if img_hash in seen_image_hashes:
                        logger.debug("Skipping duplicate image on slide %d", i)
                        continue
                    seen_image_hashes.add(img_hash)

                    image_ext = shape.image.ext
                    image_filename = f"{safe_filename}_page{i}_img{img_index}.{image_ext}"
                    image_path = os.path.join(images_dir, image_filename)
                    with open(image_path, "wb") as f:
                        f.write(image_bytes)
                    img_index += 1

                    try:
                        import io
                        img = Image.open(io.BytesIO(image_bytes))
                        ocr_text = pytesseract.image_to_string(img).strip()
                        if ocr_text:
                            slide_summaries[i] += f"\n[Diagram Text OCR: {ocr_text}]\n"
                    except Exception as e:
                        logger.warning("OCR Failed on image: %s", e)

                    summary = summarize_image(image_bytes)
                    if summary:
                        slide_summaries[i] += summary
    except Exception as e:
        logger.error("Error parsing PPTX images %s: %s", filename, e)

    # 2. Extract semantic text using Unstructured
    try:
        elements = partition_pptx(filename=file_path)

        slide_texts = {}
        for el in elements:
            page_num = el.metadata.page_number
            if not page_num:
                page_num = 1

            if page_num not in slide_texts:
                slide_texts[page_num] = []

            if el.text.strip():
                slide_texts[page_num].append(el.text.strip())

        all_pages = set(slide_texts.keys()).union(set(slide_summaries.keys()))

        for page_num in sorted(list(all_pages)):
            combined_text = ""

            if page_num in slide_texts:
                combined_text += "\n".join(slide_texts[page_num]) + "\n"

            if page_num in slide_summaries and slide_summaries[page_num]:
                combined_text += "\n" + slide_summaries[page_num] + "\n"

            combined_text = combined_text.strip()

            if combined_text:
                docs.append({
                    "text": combined_text,
                    "source": filename,
                    "page": page_num,
                    "type": "pptx"
                })
    except Exception as e:
        logger.error("Failed to parse pptx %s: %s", filename, e)

    return docs


def parse_file(file_path: str, images_dir: str = None) -> list[dict]:
    """Parse a file based on its extension. Returns list of document dicts."""
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return parse_pdf(file_path, images_dir)
    elif ext == ".docx":
        return parse_docx(file_path, images_dir)
    elif ext == ".pptx":
        return parse_pptx(file_path, images_dir)
    else:
        raise ValueError(f"Unsupported file format: {ext}")


def parse_file_streaming(file_path: str, images_dir: str = None):
    """Parse a file with progress streaming. Only PDFs support per-page streaming.
    For other formats, yields all docs at once.

    Yields: (doc_dict, current_page, total_pages)
    """
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        yield from parse_pdf_streaming(file_path, images_dir)
    else:
        docs = parse_file(file_path, images_dir)
        total = len(docs)
        for i, doc in enumerate(docs, 1):
            yield doc, i, total
